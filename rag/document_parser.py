import os
import re
from typing import List, Dict, Any

class DocumentParser:
    """
    Parses PDF, DOCX, PPTX, TXT, and Markdown files into raw text strings.
    Handles fallbacks gracefully if optional format libraries are absent.
    """

    @staticmethod
    def parse_file(file_bytes: bytes, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()

        if ext in [".txt", ".md", ".json"]:
            return file_bytes.decode("utf-8", errors="ignore")

        if ext == ".pdf":
            try:
                import pypdf
                import io
                reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                text = []
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text.append(t)
                return "\n".join(text)
            except Exception:
                return file_bytes.decode("utf-8", errors="ignore")

        if ext == ".docx":
            try:
                import docx
                import io
                doc = docx.Document(io.BytesIO(file_bytes))
                return "\n".join([p.text for p in doc.paragraphs if p.text])
            except Exception:
                return file_bytes.decode("utf-8", errors="ignore")

        if ext == ".pptx":
            try:
                import pptx
                import io
                prs = pptx.Presentation(io.BytesIO(file_bytes))
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text.append(shape.text)
                return "\n".join(text)
            except Exception:
                return file_bytes.decode("utf-8", errors="ignore")

        return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        words = re.split(r'\s+', text.strip())
        if not words or not words[0]:
            return []

        chunks = []
        i = 0
        while i < len(words):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
            i += (chunk_size - overlap)
        return chunks
